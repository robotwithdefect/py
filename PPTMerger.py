#pip install python-pptx
import collections.abc
from pptx import Presentation
import os

srcFolder = r'D:\\OneDrive\\Learning\\DS\\Sem2\\ISM\\Midsem'

for root,dirs,files in os.walk(srcFolder):
    pptFiles = [
                os.path.join(srcFolder,f) for f in files 
                if (f.endswith(".pptx") and not f.startswith("~"))
              ]

    # Open the first PowerPoint file
    prs1 = Presentation(pptFiles[0])              
    for file in pptFiles:
        if file == pptFiles[0]:
            continue
        # Open the second PowerPoint file
        prs2 = Presentation(file)
        #print(type(prs1.slides))
        # Iterate through the slides in the second PowerPoint file
        for slide in prs2.slides:
            # Add the slide to the first PowerPoint file
            sl = prs1.slides.add_slide(prs2.slide_layouts[0])
                    
    # Save the merged PowerPoint file
prs1.save('{}\\merged.pptx'.format(srcFolder))





    


